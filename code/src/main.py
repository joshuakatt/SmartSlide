import collections
import glob
import logging
import subprocess
import textrazor
from concurrent.futures import ThreadPoolExecutor
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import tkinter as tk
from tkinter import filedialog
import tkinter.scrolledtext as tkst
import speech_recognition as sr
import queue
import threading
import time
import warnings
from urllib3.exceptions import InsecureRequestWarning
warnings.simplefilter('ignore', category=InsecureRequestWarning)

from GUI import Application

class SpeechRecognizer:
    def __init__(self, pptData, message_queue):
        self.pptData = pptData
        self.message_queue = message_queue
        self.recognizer = sr.Recognizer()
        self.microphone = sr.Microphone()
        with self.microphone as source:
            self.recognizer.adjust_for_ambient_noise(source)
        self.stop_listening = self.recognizer.listen_in_background(self.microphone, self.callback)

    def callback(self, recognizer, audio):
        try:
            recognized_speech = recognizer.recognize_google(audio)
            self.message_queue.put("Recognized Speech: " + recognized_speech)
            for index, slide_text in enumerate(self.pptData):
                if check_relevance(recognized_speech, slide_text):
                    change_slide(index)
        except sr.UnknownValueError:
            self.message_queue.put("Google Speech Recognition could not understand audio")
        except sr.RequestError as e:
            self.message_queue.put(f"Could not request results from Google Speech Recognition service; {e}")

def auto_advance_slides(num_slides, message_queue):
    for i in range(num_slides):
        change_slide(i)
        time.sleep(5)

def check_relevance(speech, slide_text):
    vectorizer = TfidfVectorizer().fit_transform([speech, slide_text])
    vectors = vectorizer.toarray()
    cos_sim = cosine_similarity(vectors)
    return cos_sim[0][1] > 0.5

def change_slide(index):
    try:
        command = f'''
        tell application "Microsoft PowerPoint"
            activate
            show slide {index + 1} of active presentation
        end tell
        '''
        process = subprocess.Popen(["osascript", "-e", command], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        output, error = process.communicate()
    except Exception as e:
        message_queue.put(f"Unexpected error in change_slide: {e}")

def extract_text_from_slides(ppt_path, message_queue):
    razor_API_KEY = "9224bb6f01f35cb07f2ec8b1c7ca780f397685e30b95d517a3eaeaee"
    textrazor.api_key = razor_API_KEY
    client = textrazor.TextRazor(extractors=["entities", "topics"])

    command = f'''
    tell application "Terminal"
        activate
        do script "open -a 'Microsoft PowerPoint' '{ppt_path}'"
    end tell
    '''

    process = subprocess.Popen(["osascript", "-e", command], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    output, error = process.communicate()

    if error:
        message_queue.put(f"Error in opening presentation: {error.decode()}")

    pptData = []
    for eachfile in glob.glob(ppt_path):
        prs = Presentation(eachfile)
        message_queue.put(f"Extracting from file: {eachfile}\n{'-' * 20}")
        slide_texts = [""] * len(prs.slides)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            texts.append(run.text)
            slide_texts[i] = "".join(texts)

        message_queue.put('\n'.join(slide_texts))
        pptData.extend(slide_texts)
    return pptData

def main():
    root = tk.Tk()
    app = Application(master=root)

    message_queue = queue.Queue()

    def update_gui():
        # Update the GUI with new messages
        while not message_queue.empty():
            message = message_queue.get()
            app.update_log(message)
        root.after(100, update_gui)

    def browse_ppt_file():
        ppt_file = filedialog.askopenfilename(filetypes=[("PowerPoint Files", "*.pptx")])
        if ppt_file:
            def run_tasks():
                pptData = extract_text_from_slides(ppt_file, message_queue)
                num_slides = len(pptData)
                SpeechRecognizer(pptData, message_queue)
                auto_advance_slides(num_slides, message_queue)

            # Run tasks in a separate thread
            threading.Thread(target=run_tasks).start()

    update_gui()

    app.browse_ppt["command"] = browse_ppt_file

    root.mainloop()

if __name__ == "__main__":
    main()
