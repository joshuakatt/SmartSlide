#!/usr/bin/env python3
import collections
import glob
import logging
import subprocess
import textrazor
import threading
from concurrent.futures import ThreadPoolExecutor
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from tkinter import filedialog, Tk, Button, ttk
import speech_recognition as sr

class SpeechRecognizer:
    def __init__(self, pptData):
        self.pptData = pptData
        self.recognizer = sr.Recognizer()
        self.microphone = sr.Microphone()
        with self.microphone as source:
            self.recognizer.adjust_for_ambient_noise(source)
        self.stop_listening = self.recognizer.listen_in_background(self.microphone, self.callback)

    def callback(self, recognizer, audio):
        try:
            recognized_speech = recognizer.recognize_google(audio)
            print("Recognized Speech: " + recognized_speech)
            for index, slide_text in enumerate(self.pptData):
                if check_relevance(recognized_speech, slide_text):
                    change_slide(index)
        except sr.UnknownValueError:
            print("Google Speech Recognition could not understand audio")
        except sr.RequestError as e:
            print(f"Could not request results from Google Speech Recognition service; {e}")

def check_relevance(speech, slide_text):
    vectorizer = TfidfVectorizer().fit_transform([speech, slide_text])
    vectors = vectorizer.toarray()
    cos_sim = cosine_similarity(vectors)
    return cos_sim[0][1] > 0.5


def change_slide(index):
    try:
        command = f'''
        tell application "Microsoft PowerPoint"
            activate
            show slide {index + 1} of active presentation
        end tell
        '''
        process = subprocess.Popen(["osascript", "-e", command], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        output, error = process.communicate()
        if error:
            logging.error(f"Error in changing slide: {error.decode()}")
        else:
            logging.info(f"Successfully changed to slide: {index + 1}")
    except Exception as e:
        logging.error(f"Unexpected error in change_slide: {e}")


def extract_text_from_slides(ppt_path):
    razor_API_KEY = "your_textrazor_api_key"
    textrazor.api_key = razor_API_KEY
    client = textrazor.TextRazor(extractors=["entities", "topics"])

    command = f'''
    tell application "Terminal"
        activate
        do script "open -a 'Microsoft PowerPoint' '{ppt_path}'"
    end tell
    '''

    process = subprocess.Popen(["osascript", "-e", command], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    output, error = process.communicate()

    if error:
        print(f"Error in opening presentation: {error.decode()}")

    pptData = []
    for eachfile in glob.glob(ppt_path):
        prs = Presentation(eachfile)
        print(f"Extracting from file: {eachfile}\n{'-' * 20}")
        slide_texts = [""] * len(prs.slides)
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slide_texts[i] += run.text
        print(slide_texts)
        pptData.extend(slide_texts)
    return pptData

def additional_logic(pptData):
    logging.info("Thread additional_logic: starting")

    razor_API_KEY = "your_textrazor_api_key"
    textrazor.api_key = razor_API_KEY
    client = textrazor.TextRazor(extractors=["entities", "topics"])

    for i, slide_text in enumerate(pptData):
        try:
            response = client.analyze(slide_text)
            entities = list(response.entities())
            entities.sort(key=lambda x: x.relevance_score, reverse=True)
            seen = set()
            for entity in entities:
                if entity.id not in seen:
                    print(f"Slide {i}: {entity.wikipedia_link}")
                    seen.add(entity.id)
        except Exception as e:
            logging.error(f"Error in additional_logic: {e}")
    logging.info("Thread additional_logic: finishing")

def browse_ppt_file():
    ppt_file = filedialog.askopenfilename(filetypes=[("PowerPoint Files", "*.pptx")])
    if ppt_file:
        with ThreadPoolExecutor(max_workers=3) as executor:
            pptData_future = executor.submit(extract_text_from_slides, ppt_file)
            pptData = pptData_future.result()
            executor.submit(SpeechRecognizer, pptData)
            executor.submit(additional_logic, pptData)

if __name__ == "__main__":
    format = "%(asctime)s: %(message)s"
    logging.basicConfig(format=format, level=logging.INFO, datefmt="%H:%M:%S")

    root = Tk()
    root.title("SmartSlide")
    root.geometry('300x200')  # Set the window size

    ttk.Style().configure("TButton", padding=6, relief="flat", background="#ccc")  # Use ttk for nicer looking buttons

    btn_browse = ttk.Button(root, text="Browse PPT", command=browse_ppt_file)
    btn_browse.pack()

    root.mainloop()
